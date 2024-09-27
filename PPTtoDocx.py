#Welcome to the PPT to Docx Converting Tool code!
#Before using this code please remember to install
#*pip install python-pptx python-docx
#*pip install tkinterdnd2

#The editing tool will show a drag box in which to drop your PPT files
#And browse files to output the converted Docx documents.


import os
from tkinter import Tk, Frame, Label, filedialog, messagebox
from tkinterdnd2 import TkinterDnD, DND_FILES
from pptx import Presentation
from docx import Document


def convert_ppt_to_docx(ppt_path, output_dir):
    try:
        # Load the PowerPoint file
        ppt = Presentation(ppt_path)
        output_file = os.path.join(output_dir,
                                   os.path.basename(ppt_path).replace('.pptx', '.docx').replace('.ppt', '.docx'))

        # Create a new Word document
        doc = Document()

        # Extract text from each slide and add it to the Word document
        for slide_num, slide in enumerate(ppt.slides, start=1):
            doc.add_heading(f"Slide {slide_num}", level=1)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    doc.add_paragraph(shape.text)

        # Save the Word document
        doc.save(output_file)
        print(f"Successfully converted {ppt_path} to {output_file}")
        return output_file

    except Exception as e:
        print(f"Error converting {ppt_path}: {e}")
        return None


def on_drop(event):
    # Get the files dropped into the window
    file_paths = event.data.split('}')  # Split based on closing brace for multiple files
    file_paths = [file.strip('{} ') for file in file_paths if file]  # Clean up the file paths
    print(f"Files dropped: {file_paths}")

    # Ask for the output directory
    output_dir = filedialog.askdirectory(title="Select Output Directory")
    if not output_dir:
        print("No output directory selected.")
        return

    # Convert each PPT/PPTX to DOCX
    for file in file_paths:
        if file.endswith('.ppt') or file.endswith('.pptx'):
            output_file = convert_ppt_to_docx(file, output_dir)
            if output_file:
                print(f"Converted file saved at: {output_file}")
            else:
                print(f"Failed to convert file: {file}")
        else:
            print(f"Unsupported file type: {file}")

    messagebox.showinfo("Conversion Complete", "All files have been converted.")


def create_gui():
    # Create main window
    root = TkinterDnD.Tk()
    root.title("PPT to DOCX Converter")

    # Frame for the drag-and-drop area
    frame = Frame(root, width=400, height=200)
    frame.pack(padx=10, pady=10)

    # Label for instructions
    label = Label(frame, text="Drag and drop PPT/PPTX files here", padx=10, pady=10)
    label.pack(pady=10)

    # Register the drop target
    root.drop_target_register(DND_FILES)
    root.dnd_bind('<<Drop>>', on_drop)

    # Start the Tkinter main loop
    print("Starting the GUI...")
    root.mainloop()


if __name__ == "__main__":
    create_gui()